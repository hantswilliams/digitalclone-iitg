"""
Presentation service for generating PowerPoint presentations
"""
import os
import datetime
from flask import current_app
from pptx import Presentation
from pptx.util import Inches
from app.services.storage_service import StorageService


class PresentationService:
    """
    Service for creating and managing PowerPoint presentations
    """
    
    def __init__(self):
        """Initialize the presentation service"""
        self.storage_service = StorageService()
    
    def create_presentation(self, title, slides_data):
        """
        Create a PowerPoint presentation with videos and text
        
        Args:
            title (str): Presentation title
            slides_data (list): List of slide data, each containing:
                - text: Slide text content
                - video_url: URL to the video for the slide
                - image_url: Optional URL to the image for the slide
            
        Returns:
            dict: Result containing status and presentation URL
        """
        try:
            # Create a new presentation
            prs = Presentation()
            
            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Add content slides
            for slide_data in slides_data:
                content_slide_layout = prs.slide_layouts[5]  # Layout with title and content
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Add title if available
                if 'title' in slide_data:
                    slide.shapes.title.text = slide_data['title']
                
                # Add text
                if 'text' in slide_data:
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(2)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    text_frame.text = slide_data['text']
                
                # Note about video - PowerPoint can't directly embed videos via python-pptx
                # Instead, we add a note about the video URL
                if 'video_url' in slide_data:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = f"Video URL: {slide_data['video_url']}"
                
                # Add image if available
                if 'image_url' in slide_data:
                    # For this, we'd need to download the image first and then add it
                    # This is a placeholder - implementation would depend on how we handle image downloads
                    current_app.logger.info(f"Would add image from {slide_data['image_url']}")
            
            # Save the presentation
            timestamp = datetime.datetime.now().strftime("%Y%m%d-%H%M%S")
            filename = f"{timestamp}_{title.replace(' ', '_')}.pptx"
            file_path = os.path.join(current_app.config['TEMP_OUTPUT_FOLDER'], filename)
            prs.save(file_path)
            
            # Upload to S3
            s3_filename = self.storage_service.upload_presentation(filename, file_path)
            
            if not s3_filename:
                return {
                    'success': False,
                    'error': 'Failed to upload presentation to storage'
                }
            
            # Build URL for the uploaded presentation
            ppt_url = f"https://{current_app.config['S3_BUCKET_NAME']}.s3.amazonaws.com/{s3_filename}"
            
            # Clean up the temporary file
            if os.path.exists(file_path):
                os.remove(file_path)
            
            return {
                'success': True,
                'url': ppt_url,
                'filename': s3_filename,
                'title': title
            }
            
        except Exception as e:
            current_app.logger.error(f"Error creating presentation: {str(e)}")
            return {
                'success': False,
                'error': f'Presentation creation error: {str(e)}'
            }
    
    def create_scorm_package(self, presentation_url, title, metadata=None):
        """
        Create a SCORM package from a PowerPoint presentation
        
        Args:
            presentation_url (str): URL to the PowerPoint presentation
            title (str): Title for the SCORM package
            metadata (dict): Optional metadata for the SCORM package
            
        Returns:
            dict: Result containing status and SCORM package URL
        """
        # This is a placeholder - SCORM package creation would require:
        # 1. Downloading the presentation
        # 2. Converting to HTML5 slides
        # 3. Adding SCORM metadata and packaging
        # 4. Uploading the package
        
        current_app.logger.info(f"SCORM package creation for {presentation_url} not yet implemented")
        
        return {
            'success': False,
            'error': 'SCORM package creation not yet implemented',
            'presentation_url': presentation_url
        }
