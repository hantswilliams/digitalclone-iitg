from flask import Flask, render_template, request, Response, stream_with_context, jsonify, url_for, session, redirect
from werkzeug.utils import secure_filename
from db_model.db import Text, Photo, Audio, Video, Powerpoint, Project, ProjectTextAssociation
from db_model.db import init_db, get_session, get_projects, get_project_data_associations, update_project_data_associations
from db_model.db import create_update_user
from utils.s3 import s3, s3_upload_video
from utils.sadtalker import create_video_job
import json
import os
import requests
import time
import uuid
import datetime
from pptx import Presentation
from pptx.util import Inches
import subprocess
import urllib
import html

from authlib.integrations.flask_client import OAuth
from authlib.common.security import generate_token

app = Flask(__name__)
app.secret_key = 'super secret key'
oauth = OAuth(app)


init_db() # Initialize the database

s3_client = s3 # Initialize the s3 client

userid = os.getenv("PLAYHT_USERID")
apikey = os.getenv("PLAYHT_SECRET")
GOOGLE_CLIENT_ID = os.getenv('GOOGLE_OAUTH_CLIENT_ID')
GOOGLE_CLIENT_SECRET = os.getenv('GOOGLE_OAUTH_CLIENT_SECRET')
BUCKET_NAME = 'iitg-mvp' # Replace with your bucket name

# Dictionary to store job_id to Gradio job object mapping
current_jobs = {}

# Dictionary to store currently job details
current_job_details = {}

ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg'}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    if 'user' not in session:
        return redirect('/google/')
    else:
        print('User logged in: ', session['user'])
        return render_template('index.html', user=session['user'])

@app.route('/google/')
def google():
    CONF_URL = 'https://accounts.google.com/.well-known/openid-configuration'
    oauth.register(
        name='google',
        client_id=GOOGLE_CLIENT_ID,
        client_secret=GOOGLE_CLIENT_SECRET,
        server_metadata_url=CONF_URL,
        client_kwargs={
            'scope': 'openid email profile'
        }
    )
    # Redirect to google_auth function note, if running locally on a non-google shell, do not need to override redirect_uri and can just use url_for as below
    redirect_uri = url_for('google_auth', _external=True)
    print('REDIRECT URL: ', redirect_uri)
    session['nonce'] = generate_token()
    ##, note: if running in google shell, need to override redirect_uri to the external web address of the shell, e.g.,
    # redirect_uri = 'https://5000-cs-213132341638-default.cs-us-east1-pkhd.cloudshell.dev/google/auth/'
    return oauth.google.authorize_redirect(redirect_uri, nonce=session['nonce'])

@app.route('/google/auth/')
def google_auth():
    token = oauth.google.authorize_access_token()
    user = oauth.google.parse_id_token(token, nonce=session['nonce'])
    session['user'] = user
    print(" Google User Information: ", session['user'])
    create_update_user(user)
    print(" Google User ", user)
    return redirect('/')

@app.route('/logout')
def logout():
    session.pop('user', None)
    return redirect('/')

@app.route('/view1')
def view1():
    return render_template('view1/view1.html')

@app.route('/view1b', methods=['GET', 'POST'])
def view1b():

    ## get images from server
    session = get_session() # Get a new session
    photo_data = session.query(Photo).all()
    session.close() # Don't forget to close the session

    print('photo_data:', photo_data)

    return render_template('view1b/view1b.html', photo_data=photo_data)

@app.route('/view1/audio/table', methods=['GET', 'POST'])
def audio_table():
    row_id = request.form['row_id']
    session = get_session()
    # Query to join DataEntry and Audio
    # query = session.query(DataEntry, Audio).join(Audio, DataEntry.id == Audio.data_entry_id).filter(DataEntry.id == row_id)
    query = session.query(Text, Audio).join(Audio, Text.id == Audio.text_id).filter(Text.id == row_id)
    data_dict = []
    # Execute the query and process results
    for text, audio in query.all():
        # add to data_dict
        data_dict.append({
            'text_content': text.text_content,
            'audio_url': audio.audio_url,
            'voice': audio.voice,
            'audio_text': audio.audio_text,
        })
        # print(f"Text Content: {data_entry.text_content}, Audio URL: {audio.audio_url}")
    print('data_dict:', data_dict)
    session.close()
    return render_template(
        'view1/audio_table_modal.html', 
        data_dict=data_dict,
    )

@app.route('/view1/form', methods=['GET'])
def view_form_audio():
    ## get projects from database
    projects = get_projects()
    return render_template(
        'view1/create_text_modal.html',
        projects = projects
        )

@app.route('/view1/image_upload_form', methods=['GET'])
def image_upload_form():
    return render_template('view1/create_image_modal.html')

@app.route('/view1/text/upload', methods=['POST'])
def submit_text_upload():
    text_input = request.form['text_input']
    return_output = request.args.get('return_type')
    print('text output type:', return_output)
    project_id_list = request.form.getlist('project_id[]')
    print('text inputted:', text_input)
    print('project_id:', project_id_list)
    session = get_session()
    data_entry = Text(
        user_id='mvp_001',
        text_content=text_input, 
    )
    try:
        session.add(data_entry)
        session.commit()
    except:
        session.rollback()
        raise

    # if project_id is not an empty list:
    if project_id_list:
        session = get_session()         ## reopen session to get the row_id
        text_id = session.query(Text).filter(Text.text_content == text_input).first().id
        print('text_id:', text_id)
        for project_id in project_id_list:
            project_text_association = ProjectTextAssociation(
                project_id=project_id,
                text_id=text_id
            )
            try:
                session.add(project_text_association)
                session.commit()
            except:
                session.rollback()
                raise
    
    session.close()
    
    if return_output == 'html':
        print('data entry added to database')
        data = {'status': '200', 
                'message': 'Text successfully added',
                'text_input': text_input}
        formatted_data = json.dumps(data, indent=4).lstrip()
        return render_template('view1/submitted_modal.html', data=formatted_data)

    else:
        print('data entry added to database')
        session = get_session()         ## reopen session to get the row_id
        text_id = session.query(Text).filter(Text.text_content == text_input).first().id
        session.close()
        data = {'status': '200', 
                'message': 'Text successfully added',
                'text_row_id': text_id}
        formatted_data = json.dumps(data, indent=4).lstrip()
        return formatted_data

    # print('data entry added to database')
    # data = {'status': '200', 
    #         'message': 'Text successfully added',
    #         'text_input': text_input}
    # formatted_data = json.dumps(data, indent=4).lstrip()
    # return render_template('view1/submitted_modal.html', data=formatted_data)

@app.route('/view1/text/edit', methods=['GET', 'POST'])
def edit_text():
    row_id = request.form['row_id']
    print('row_id:', row_id)

    session = get_session()
    data_entry = session.query(Text).filter(Text.id == row_id).first()
    session.close()
    data_entry.text_content = data_entry.text_content.strip()
    print('Data entry: ', data_entry)
    
    ## get list of associated projects
    project_list_associated, project_list_non_associated = get_project_data_associations('text', row_id)

    return render_template(
        'view1/edit_text_modal.html',
        data=data_entry,
        project_list_associated=project_list_associated,
        project_list_non_associated=project_list_non_associated
        )

@app.route('/view1/text/update', methods=['POST'])
def update_text():
    row_id = request.form['row_id']
    text_input = request.form['text_input']
    associated_projects = request.form.getlist('project_id_associated')
    non_associated_projects = request.form.getlist('project_id_non_associated')
    projects_to_keep = associated_projects + non_associated_projects
    session = get_session()
    data_entry = session.query(Text).filter(Text.id == row_id).first()
    data_entry.text_content = text_input
    try:
        session.commit()
    except:
        session.rollback()
        raise

    # update associations via update_project_data_associations
    update_project_data_associations('text', row_id, projects_to_keep)

    print('data entry updated in database')
    data = {'status': '200', 
            'message': 'Text successfully updated',
            'text_input': text_input}
    formatted_data = json.dumps(data, indent=4).lstrip()
    return render_template('view1/submitted_modal.html', data=formatted_data)

@app.route('/view1/image/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        data = {'status': '400', 'message': 'No file part'}
        return render_template('view1/submitted_modal.html', data=data)

    file = request.files['file']
    if file.filename == '':
        data = {'status': '400', 'message': 'No selected file'}
        return render_template('view1/submitted_modal.html', data=data)

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        s3_client.upload_fileobj(file, BUCKET_NAME, filename)
        # Construct the file URL
        file_url = f"https://{BUCKET_NAME}.s3.amazonaws.com/{filename}"
        print('CALCULATED FILE URL:', file_url)
        # Save the file URL to the database along with meta data
        session = get_session() # Get a new session
        data_entry = Photo(
            user_id='mvp_001',
            photo_description='Photo from User X', 
            photo_url=file_url
        )
        try:
            session.add(data_entry)
            session.commit()
        except:
            session.rollback()
            raise
        session.close()
        print('file uploaded successfully to database and s3')
        data = {
            'status': '200', 
            'message': 
            'File successfully uploaded',
            'file_url': file_url
            }
        formatted_data = json.dumps(data, indent=4).lstrip()
        return render_template('view1/submitted_modal.html', data=formatted_data)
    
    else:
        data = {'status': '400', 'message': 'File type not allowed'}
        return render_template('view1/submitted_modal.html', data=data)

@app.route('/view1/text/delete', methods=['POST'])
def delete_text():
    row_id = request.form['row_id']
    print('row_id:', row_id)
    session = get_session()
    try:
        ## delete from DataEntry
        session.query(Text).filter(Text.id == row_id).delete()
        ## delete any associated audio
        session.query(Audio).filter(Audio.text_id == row_id).delete()
        ## delete associated project_text_association
        session.query(ProjectTextAssociation).filter(ProjectTextAssociation.text_id == row_id).delete()
        ## commit changes
        session.commit()
    except:
        session.rollback()
        raise
    session.close()
    print({
            'status': '200', 
            'message': 'Text successfully deleted',
            'row_id': row_id
        })
    data = {
            'status': '200', 
            'message': 'Text successfully deleted',
            'row_id': row_id
        }
    formatted_data = json.dumps(data, indent=4).lstrip()
    return render_template('view1/submitted_modal.html', data=formatted_data)

@app.route('/view1/image/delete', methods=['POST'])
def delete_file():
    photo_url = request.form['photo_url']
    file_name = photo_url.split('/')[-1]
    print('photo_url to delete:', file_name)
    try:
        s3_client.delete_object(Bucket=BUCKET_NAME, Key=file_name)
        print(f"Deleted {file_name} from {BUCKET_NAME}")
    except Exception as e:
        print(f"Error: {e}")
        return None
    session = get_session()
    try:
        session.query(Photo).filter(Photo.photo_url == photo_url).delete()
        session.commit()
    except:
        session.rollback()
        raise
    session.close()
    print({
            'status': '200', 
            'message': 'File successfully deleted',
            'photo_url': photo_url
        })
    data = {
            'status': '200', 
            'message': 'File successfully deleted',
            'photo_url': photo_url
        }
    formatted_data = json.dumps(data, indent=4).lstrip()
    return render_template('view1/submitted_modal.html', data=formatted_data)

@app.route('/view1/image/edit', methods=['GET', 'POST'])
def image_edit():
    row_id = request.form['row_id']
    print('row_id:', row_id)

    session = get_session()
    data_entry = session.query(Photo).filter(Photo.id == row_id).first()
    session.close()
    data_entry.photo_description = data_entry.photo_description.strip()
    print('Data entry: ', data_entry)
    
    ## get list of associated projects
    project_list_associated, project_list_non_associated = get_project_data_associations('photo', row_id)

    return render_template(
        'view1/edit_image_modal.html',
        data=data_entry,
        project_list_associated=project_list_associated,
        project_list_non_associated=project_list_non_associated
        )

@app.route('/view1/image/update', methods=['POST'])
def image_update():
    row_id = request.form['row_id']
    image_description = request.form['image_description']
    associated_projects = request.form.getlist('project_id_associated')
    non_associated_projects = request.form.getlist('project_id_non_associated')
    projects_to_keep = associated_projects + non_associated_projects
    session = get_session()
    data_entry = session.query(Photo).filter(Photo.id == row_id).first()
    data_entry.photo_description = image_description
    try:
        session.commit()
    except:
        session.rollback()
        raise

    # update associations via update_project_data_associations
    print('projects_to_keep:', projects_to_keep)
    update_project_data_associations('photo', row_id, projects_to_keep)

    print('data entry updated in database')
    data = {'status': '200', 
            'message': 'Text successfully updated',
            'image_description': image_description}
    formatted_data = json.dumps(data, indent=4).lstrip()
    return render_template('view1/submitted_modal.html', data=formatted_data)

@app.route('/view1/project/create', methods=['GET', 'POST'])
def create_project():
    return render_template('view1/create_project_modal.html')

@app.route('/view1/project/create/submit', methods=['POST'])
def create_project_submit():
    ## from form get project_name and project_description
    project_name = request.form['project_name']
    project_description = request.form['project_description']
    print(f'project_name: {project_name}, project_description: {project_description}')
    ## save to database in Project table
    session = get_session()
    try:
        new_project = Project(
            project_name=project_name,
            project_description=project_description
        )
        session.add(new_project)
        session.commit()
        session.close()
    except:
        session.rollback()
        raise
    print('project saved to database')
    data = {'status': '200', 'message': 'Project successfully created'}
    formatted_data = json.dumps(data, indent=4).lstrip()
    return render_template('view1/submitted_modal.html', data=formatted_data)

@app.route('/view1/project/delete', methods=['POST'])
def delete_project():
    project_id = request.form['row_id']
    print('project_id:', project_id)
    session = get_session()
    try:
        ## delete from Project
        session.query(Project).filter(Project.id == project_id).delete()
        ## delete any associated audio
        # session.query(Text).filter(Text.project_id == project_id).delete()
        # session.query(Audio).filter(Audio.project_id == project_id).delete()
        # session.query(Photo).filter(Photo.project_id == project_id).delete()
        # session.query(Video).filter(Video.project_id == project_id).delete()
        # session.query(Powerpoint).filter(Powerpoint.project_id == project_id).delete()
        session.commit()
    except:
        session.rollback()
        raise
    session.close()
    print({
            'status': '200', 
            'message': 'Project successfully deleted',
            'project_id': project_id
        })
    data = {
            'status': '200', 
            'message': 'Project successfully deleted',
            'project_id': project_id
        }
    formatted_data = json.dumps(data, indent=4).lstrip()
    return render_template('view1/submitted_modal.html', data=formatted_data)

@app.route('/view1/project/edit', methods=['GET', 'POST'])
def edit_project():
    project_id = request.form['row_id']
    print('project_id:', project_id)
    session = get_session()
    project = session.query(Project).filter(Project.id == project_id).first()
    print('Project: ', project)
    return render_template(
        'view1/edit_project_modal.html',
        data=project
        )

@app.route('/view1/project/update', methods=['POST'])
def update_project():
    project_id = request.form['row_id']
    project_name = request.form['project_name']
    project_description = request.form['project_description']
    print('project_id:', project_id)
    print('project_name:', project_name)
    print('project_description:', project_description)
    session = get_session()
    project = session.query(Project).filter(Project.id == project_id).first()
    project.project_name = project_name
    project.project_description = project_description
    try:
        session.commit()
    except:
        session.rollback()
        raise
    session.close()
    print('project updated in database')
    data = {'status': '200', 
            'message': 'Project successfully updated',
            'project_name': project_name}
    formatted_data = json.dumps(data, indent=4).lstrip()
    return render_template('view1/submitted_modal.html', data=formatted_data)

@app.route('/view1/project/modal', methods=['GET', 'POST'])
def project_modal():
    ## get projects from database
    projects = get_projects()
    print('projects:', projects)
    ## loop through and print each project
    for project in projects:
        print(f"Project ID: {project.id}, Project Name: {project.project_name}, Project Description: {project.project_description}")
    return render_template('view1/project_table_modal.html', data=projects)

@app.route('/view2', methods=['GET', 'POST'])
def view2():

    session = get_session()  # Get a new session

    # text_data = [(item.text_content, item.data_url) for item in all_data if item.data_type == 'Text']
    # photo_data = [(item.text_content, item.data_url) for item in all_data if item.data_type == 'Photo']
    
    text_data = session.query(Text).all()
    photo_data = session.query(Photo).all()

    for i in text_data:
        print('XXX Text data: ', i.text_content)

    for i in photo_data:
        print('XXX Photo Data: ', i.photo_url)

    # Query to join DataEntry and Audio
    # query = session.query(DataEntry, Audio).join(Audio, DataEntry.id == Audio.data_entry_id)
    query = session.query(Text, Audio).join(Audio, Text.id == Audio.text_id)

    print('Combined QUERY data: ', query.all())

    data_dict = []
    # Execute the query and process results
    for text, audio in query.all():
        # add to data_dict
        data_dict.append({
            'text_content': text.text_content,
            'audio_url': audio.audio_url,
            'voice': audio.voice,
            'audio_text': audio.audio_text,
        })
        # print(f"Text Content: {data_entry.text_content}, Audio URL: {audio.audio_url}")
    print('data_dict:', data_dict)
    session.close()  # Don't forget to close the session

    return render_template(
        'view2/view2.html', 
        data_dict=data_dict,
        text_data=text_data, 
        photo_data=photo_data
    )

@app.route('/view2/submit', methods=['POST'])
def submit():
    print('DATA FORM RECEIVED', request.form)
    # data = json.dumps(request.form)

    print('Current Jobs:', current_jobs)
    print('Current Form Text/Audio:', request.form['audio_selection'])
    print('Current Form Photo:', request.form['photo_selection'])
    
    # audio_selection = request.form['audio_selection']
    audio_text_values = request.form['audio_selection']
    audio_selection, text_selection = audio_text_values.split('||', 1)
    photo_selection = request.form['photo_selection']
    return_type = request.form['return_type']

    print('audio_selection:', audio_selection)
    print('photo_selection:', photo_selection)

    try:
        job = create_video_job(photo_selection, audio_selection)
        print('Initial State at Launch:', job.status())
    except Exception as e:
        print(f"Error: {e}")
        return None

    job_id = 'job_' + str(uuid.uuid4())[:4]
    current_jobs[job_id] = job

    ## get row ID from db audio for the audio_selection
    session = get_session()
    audio_row = session.query(Audio).filter(Audio.audio_url == audio_selection).first()
    # audio_row_id = audio_row.id
    text_row_id = audio_row.text_id
    # print('audio_row_id:', audio_row_id)

    ## get row ID from db photo for the photo_selection
    photo_row = session.query(Photo).filter(Photo.photo_url == photo_selection).first()
    photo_row_id = photo_row.id
    print('photo_row_id:', photo_row_id)

    current_job_details[job_id] = {
        'job_id': job_id,
        'audio_selection': audio_selection,
        # 'audio_row_id': audio_row_id,
        'text_selection': text_selection,
        'text_row_id': text_row_id,
        'photo_selection': photo_selection,
        'photo_row_id': photo_row_id,
        'start_time': datetime.datetime.now(),
    }

    results = {
        'job_id': job_id,
        'job_status': job.status().code,
    }
    
    if return_type == 'html':
        return render_template('view2/submitted_modal.html', data=results)
    else:
        return job_id

@app.route('/view2/submit/progress')
def submit_progress():
    ## get the job_id from the request
    print('Current Jobs:', current_jobs)
    job_id = request.args.get('job_id')
    passed_job = job_id
    ## get the job from the current_jobs dictionary
    passed_job_object = current_jobs[passed_job]
    ## get the job details from the current_job_details dictionary
    passed_job_details = current_job_details[passed_job]
    print('passed_job_details from serverside:', passed_job_details)
    print('Before generate status():', passed_job_object.status())
    
    
    def generate():
        while True:
            job_status_code = str(passed_job_object.status().code)
            
            print(f"Current job_status: {job_status_code}")

            if job_status_code in ["Status.STARTING", "Status.PROCESSING"]:
                yield f"data: {job_status_code}\n\n"
                time.sleep(1)

            elif job_status_code == "Status.FINISHED":

                print('FINISHED job_status:', job_status_code)
                yield f"data: {job_status_code}\n\n"
                
                time.sleep(5) # made this slower to allow for video download from remote server
                ## if on a faster machine, get make this quicker - currently out of country so need to do this

                ## get working directory
                working_dir = os.getcwd()
                print('WORKING DIR: ', working_dir)

                ## get the new folder name from the temp_outputs folder
                folder = os.listdir("./temp_outputs")

                ## get only the folders that exist in the temp_outputs folder
                folders = [f for f in folder if os.path.isdir(os.path.join("./temp_outputs", f))]
                print('FOLDERS FOUND IN TEMP:', folders)

                ###############################################################
                ###############################################################
                #####################FOR DEALING WITH SLOW INTERNET############
                ###############################################################
                ###############################################################

                start_time = time.time()  # Capture the start time
                timeout = 60  # Set the timeout period to 60 seconds

                while True:
                    # List only directories in the specified path
                    folders = [f for f in os.listdir("./temp_outputs") if os.path.isdir(os.path.join("./temp_outputs", f))]
                    print('FOLDERS FOUND IN TEMP:', folders)
                    
                    # Check if the folder list is not empty or if the timeout has been reached
                    if folders or (time.time() - start_time) > timeout:
                        break  # Exit the loop if folders are found or timeout has been reached
                    
                    time.sleep(5)  # Wait for 5 seconds before trying again

                # Check the reason for the loop exit and act accordingly
                if not folders:
                    print('Error: No folders found within the 60 seconds time frame.')
                    # Handle the error condition here
                else:
                    print('Folders have been found.')
                    # Continue with the rest of your code
                    
                ###############################################################
                ###############################################################
                #####################FOR DEALING WITH SLOW INTERNET############
                ###############################################################
                ###############################################################

                ## get the folder name
                folder_name = folders[0]
                # print('FOLDER NAME: ', folder_name)
                
                ## get the .mp4 file from the folder
                video_file = os.listdir("./temp_outputs/" + folder_name)[0]
                print('FOUND NEW VID FILE: ', video_file)
                
                ## upload the video file to S3
                try:
                    video_url = s3_upload_video(video_file, "temp_outputs/" + folder_name + "/" + video_file)
                    print('Succesfully uploaded video from: ', video_url)
                    ## then delete the folder
                    os.system("rm -rf temp_outputs/" + folder_name)
                    print(f"Deleted folder {folder_name}")
                
                except Exception as e:
                    print(f"Error: {e}")
                    return None
                
                ## save to database in Video table, providing the row_id as the data_entry_id
                session = get_session()
                
                try:
                    new_video = Video(
                        job_id=job_id,
                        audio_selection=passed_job_details['audio_selection'],
                        # audio_row_id=passed_job_details['audio_row_id'],
                        text_selection=passed_job_details['text_selection'],
                        text_row_id=passed_job_details['text_row_id'],
                        photo_selection=passed_job_details['photo_selection'],
                        photo_row_id=passed_job_details['photo_row_id'],
                        video_url=f'https://{BUCKET_NAME}.s3.amazonaws.com/{video_url}'
                    )
                    session.add(new_video)
                    session.commit()
                    session.close()
                    ## yield the video_url
                    s3_url=f'https://{BUCKET_NAME}.s3.amazonaws.com/{video_url}'
                    yield f"data: {s3_url}\n\n"
                    ## delete the job from the current_jobs dictionary
                    del current_jobs[passed_job]
            
                except:
                    session.rollback()
                    raise

                # Break the loop after sending the final status
                break

            else:
                # Handle any other statuses if necessary
                yield f"data: {job_status_code}"
                break  # Or continue, based on your requirements

            time.sleep(1)  # Sleep at the end of the loop

        # After loop ends
        yield "data: Done\n\n"    
        
    # def generate():
    #     while str(passed_job_object.status().code) == "Status.STARTING":
    #         job_status = str(passed_job_object.status().code)
    #         print('STARTING job_status:', job_status)
    #         yield f"data: {job_status}\n\n"
    #         time.sleep(1)
        
    #     while str(passed_job_object.status().code) == "Status.PROCESSING":
    #         job_status = str(passed_job_object.status().code)
    #         print('PROCESSING job_status:', job_status)
    #         yield f"data: {job_status}\n\n"
    #         time.sleep(1)
        
    #     while str(passed_job_object.status().code) == "Status.FINISHED":
    #         job_status = str(passed_job_object.status().code)
    #         print('FINISHED job_status:', job_status)
    #         yield f"data: {job_status}\n\n"
            
    #         time.sleep(5) # made this slower to allow for video download from remote server
    #         ## if on a faster machine, get make this quicker - currently out of country so need to do this

    #         ## get working directory
    #         working_dir = os.getcwd()
    #         print('WORKING DIR: ', working_dir)

    #         ## get the new folder name from the temp_outputs folder
    #         folder = os.listdir("./temp_outputs")

    #         ## get only the folders that exist in the temp_outputs folder
    #         folders = [f for f in folder if os.path.isdir(os.path.join("./temp_outputs", f))]
    #         print('FOLDERS FOUND IN TEMP:', folders)

    #         ###############################################################
    #         ###############################################################
    #         #####################FOR DEALING WITH SLOW INTERNET############
    #         ###############################################################
    #         ###############################################################

    #         start_time = time.time()  # Capture the start time
    #         timeout = 60  # Set the timeout period to 60 seconds

    #         while True:
    #             # List only directories in the specified path
    #             folders = [f for f in os.listdir("./temp_outputs") if os.path.isdir(os.path.join("./temp_outputs", f))]
    #             print('FOLDERS FOUND IN TEMP:', folders)
                
    #             # Check if the folder list is not empty or if the timeout has been reached
    #             if folders or (time.time() - start_time) > timeout:
    #                 break  # Exit the loop if folders are found or timeout has been reached
                
    #             time.sleep(5)  # Wait for 5 seconds before trying again

    #         # Check the reason for the loop exit and act accordingly
    #         if not folders:
    #             print('Error: No folders found within the 60 seconds time frame.')
    #             # Handle the error condition here
    #         else:
    #             print('Folders have been found.')
    #             # Continue with the rest of your code
                
    #         ###############################################################
    #         ###############################################################
    #         #####################FOR DEALING WITH SLOW INTERNET############
    #         ###############################################################
    #         ###############################################################

    #         ## get the folder name
    #         folder_name = folders[0]
    #         # print('FOLDER NAME: ', folder_name)
            
    #         ## get the .mp4 file from the folder
    #         video_file = os.listdir("./temp_outputs/" + folder_name)[0]
    #         print('FOUND NEW VID FILE: ', video_file)
            
    #         ## upload the video file to S3
    #         try:
    #             video_url = s3_upload_video(video_file, "temp_outputs/" + folder_name + "/" + video_file)
    #             print('Succesfully uploaded video from: ', video_url)
    #             ## then delete the folder
    #             os.system("rm -rf temp_outputs/" + folder_name)
    #             print(f"Deleted folder {folder_name}")
            
    #         except Exception as e:
    #             print(f"Error: {e}")
    #             return None
            
    #         ## save to database in Video table, providing the row_id as the data_entry_id
    #         session = get_session()
    #         try:
    #             new_video = Video(
    #                 job_id=job_id,
    #                 audio_selection=passed_job_details['audio_selection'],
    #                 # audio_row_id=passed_job_details['audio_row_id'],
    #                 text_selection=passed_job_details['text_selection'],
    #                 text_row_id=passed_job_details['text_row_id'],
    #                 photo_selection=passed_job_details['photo_selection'],
    #                 photo_row_id=passed_job_details['photo_row_id'],
    #                 video_url=f'https://{BUCKET_NAME}.s3.amazonaws.com/{video_url}'
    #             )
    #             session.add(new_video)
    #             session.commit()
    #             session.close()
    #             ## yield the video_url
    #             s3_url=f'https://{BUCKET_NAME}.s3.amazonaws.com/{video_url}'
    #             yield f"data: {s3_url}\n\n"
    #             ## delete the job from the current_jobs dictionary
    #             del current_jobs[passed_job]
            
    #         except:
    #             session.rollback()
    #             raise            
    #     else:
    #         job_status = str(passed_job_object.status().code)
    #         print('OTHER job_status:', job_status)
    #         yield f"data: {job_status}"
        
    #     yield f"data: {job_status}"

    return Response(stream_with_context(generate()), mimetype='text/event-stream')

@app.route('/view2b', methods=['GET', 'POST'])
def view2b():
    session = get_session()  # Get a new session
    photo_data = session.query(Photo).all()
    session.close()
    return render_template(
        'view2/view2b.html',
        photo_data=photo_data
    )

@app.route("/view2b/submit/part1", methods=["POST"])
def submit2b():
    # Part 0 - Get the data from the form for multiple slides
    slide_titles = request.form.getlist('slide_title[]')
    slide_bodies = request.form.getlist('slide_body[]')
    slide_audio_texts = request.form.getlist('slide_audio_text[]')
    voice_selections = request.form.getlist('voice[]')
    image_selections = request.form.getlist('image_select[]')

    # Process each slide's data
    slides_data = []
    for i in range(len(slide_titles)):
        slide_data = {
            'title': slide_titles[i],
            'body': slide_bodies[i].replace('\r\n', '\n').replace('\r', '\n'),
            'audio_text': slide_audio_texts[i],
            'voice': voice_selections[i],
            'image': image_selections[i]
        }
        slides_data.append(slide_data)

        # Part 1 - Save text for each slide (example shown for the first slide)
        if i == 0:  # As an example, processing only the first slide
            url = "http://localhost:5005/view1/text/upload"
            data = {"text_input": slide_audio_texts[i]}
            headers = {'Content-Type': 'application/x-www-form-urlencoded'}
            response = requests.post(url, data=data, headers=headers)
            text_row_id = response.json().get('text_row_id', None)
        else:
            text_row_id = None  # Placeholder for other slides processing

        # Print slide data for debugging
        print(f"FROM SERVER: Slide {i+1}: {slide_data}")

    ## Part 2 - Create audio file (and other processing as needed)
    return render_template(
        'view2/submitted_modal_2b.html',
        slides_data=slides_data,
        text_row_id=text_row_id  # Example: sending text_row_id of the first slide
    )

# @app.route("/view2b/submit/part1", methods=["POST"])
# def submit2b():

#     ## Part 0 - Get the data from the form
#     slide_title = request.form['slide_title']

#     slide_body = request.form['slide_body'].replace('\r\n', '\n').replace('\r', '\n')
#     slide_body = json.dumps(slide_body)

#     # slide_body = request.form['slide_body'].replace('\r\n', '\n').replace('\r', '\n')
#     # print('slide_body_from_server: ', slide_body)

#     # slide_body = request.form['slide_body'].replace('\r\n', '\n').replace('\r', '\n')
#     # slide_body = json.dumps(slide_body)

#     slide_audio_text = request.form['slide_audio_text']
#     voice_selection = request.form['voice']
#     image_selection = request.form['image_select']
#     print(f"FROM SERVER: slide_title: {slide_title}, slide_body: {slide_body}, slide_audio_text: {slide_audio_text}, voice_selection: {voice_selection}, image_selection: {image_selection}")
    
#     ## Part 1 - Save text
#     ### send data to /view1/text/upload endpoint 
#     url = "http://localhost:5005/view1/text/upload"
#     data = {"text_input": slide_audio_text}
#     headers={'Content-Type': 'application/x-www-form-urlencoded'}
#     response = requests.post(url, data=data, headers=headers)
#     response.json()
#     text_row_id = response.json()['text_row_id']

#     ## Part 2 - Create audio file    
#     return render_template(
#         'view2/submitted_modal_2b.html',
#         text=slide_audio_text,
#         voice=request.form['voice'],
#         text_row_id=text_row_id,
#         image_selection=image_selection,
#         slide_title=slide_title,
#         slide_body=slide_body
#     )

@app.route('/view3')
def view3():
    ## get videos from video table
    session = get_session()
    ## exclude rows with no video_url
    video_data = session.query(Video).all()
    session.close()
    print('video_data:', video_data)
    # s3_video_urls = s3_get_mp4_files()
    # print('s3_video_urls:', s3_video_urls)
    return render_template(
        'view3/view3.html', 
        data=video_data
    )

@app.route('/view3/video/delete', methods=['POST'])
def delete_video():
    ## from form get row id
    row_id = request.form['row_id']
    print('row_id from req args delete:', row_id)
    session = get_session()
    try:
        ## delete from Video
        session.query(Video).filter(Video.id == row_id).delete()
        session.commit()
    except:
        session.rollback()
        raise
    session.close()
    print({
            'status': '200', 
            'message': 'Video successfully deleted',
            'row_id': row_id
        })
    data = {
            'status': '200', 
            'message': 'Video successfully deleted',
            'row_id': row_id
        }
    formatted_data = json.dumps(data, indent=4).lstrip()
    return render_template('view1/submitted_modal.html', data=formatted_data)

@app.route('/data/view1', methods=['GET'])
def data():
    if request.headers.get('HX-Request'):

        session = get_session() # Get a new session
        # data = session.query(DataEntry).order_by(func.random()).limit(5).all()
        # data = session.query(DataEntry).all()
        text_data = session.query(Text).all()
        photo_data = session.query(Photo).all()
        session.close() # Don't forget to close the session

        ## add in data_type to each row
        for row in text_data:
            row.data_type = 'Text'
        for row in photo_data:
            row.data_type = 'Photo'

        # combine text and photo data
        data = text_data + photo_data

        # ## get s3 image and sound urls
        # s3_image_urls, s3_sound_urls = s3_get_image_sounds()
        # print('s3_image_urls:', s3_image_urls)
        # print('s3_sound_urls:', s3_sound_urls)

        return render_template('view1/data_table.html', data=data)
    else:
        return render_template('view1.html', data=data)

@app.route('/audio/create', methods=['GET', 'POST'])
def stream_page():
    row_id = request.form['row_id']
    print('row_id:', row_id)
    session = get_session()
    data_entry = session.query(Text).filter(Text.id == row_id).first()
    print('Data entry: ', data_entry)
    return render_template(
        'playht/create_audio_modal.html',
        data=data_entry
    )

@app.route('/audio/stream-data', methods=['GET', 'POST'])
def stream():

    text = request.args.get('text', 'Default text')  # Fallback to default text if not provided
    voice = request.args.get('voice', 'Default voice')  # Fallback to default voice if not provided
    stream = request.args.get('stream', 'False')  # Fallback to default voice if not provided

    print(f'CREATING AUDIO: text: {text}, voice: {voice}, stream: {stream}')

    voices = {
        "male_matt": "s3://voice-cloning-zero-shot/09b5c0cc-a8f4-4450-aaab-3657b9965d0b/podcaster/manifest.json",
        "female_nichole": "s3://voice-cloning-zero-shot/7c38b588-14e8-42b9-bacd-e03d1d673c3c/nicole/manifest.json",
        "male_hants": "s3://voice-cloning-zero-shot/fdbb6c51-284f-4ffb-8820-56d3f5513862/hants-v1/manifest.json",
    }

    if stream == 'True':
        def event_stream():
            url = "https://api.play.ht/api/v2/tts"
            payload = {
                "text": f"{str(text)}",      # "Haunts",
                "voice": f"{voices[voice]}",     # "s3://voice-cloning-zero-shot/09b5c0cc-a8f4-4450-aaab-3657b9965d0b/podcaster/manifest.json",
                "output_format": "wav",
                "voice_engine": "PlayHT2.0"
            }
            headers = {
                "accept": "text/event-stream",
                "content-type": "application/json",
                "Authorization": "Bearer " + apikey,
                "X-USER-ID": userid
            }
            with requests.post(url, stream=True, headers=headers, json=payload) as response:
                for line in response.iter_lines():
                    if line:
                        decoded_line = line.decode('utf-8')
                        yield f"data: {decoded_line}\n\n"

        return Response(
            stream_with_context(event_stream()),
            content_type='text/event-stream'
            )
    
    elif stream == 'False':
        url = "https://api.play.ht/api/v2/tts"
        payload = {
            "text": f"{str(text)}",      # "Haunts",
            "voice": f"{voices[voice]}",     # "s3://voice-cloning-zero-shot/09b5c0cc-a8f4-4450-aaab-3657b9965d0b/podcaster/manifest.json",
            "output_format": "wav",
            "voice_engine": "PlayHT2.0"
        }
        headers = {
            "accept": "application/json",
            "content-type": "application/json",
            "Authorization": "Bearer " + apikey,
            "X-USER-ID": userid
        }
        response = requests.post(url, headers=headers, json=payload)
        print('playHT job ID response:', response)

        return response.json()

@app.route('/audio/stream-data/response', methods=['POST'])
def stream_response():
    body = request.get_json()
    row_id = body['row_id'] 
    audio_url = body['playht']['url']
    voice = body['voice']
    text_content = body['text']
    print(f'FROM SERVER: row_id: {row_id}, audio_url: {audio_url}, voice: {voice}')
    session = get_session()
    try:
        new_audio = Audio(
            audio_url=audio_url, 
            voice=voice, 
            audio_text=f"{str(text_content)}",
            # data_entry_id=row_id
            text_id=row_id
        )
        session.add(new_audio)
        session.commit()
        session.close()
    except:
        session.rollback()
        raise
    print('audio saved to database')
    return 'success - audio saved to database'

@app.route('/view1/delete/db', methods=['GET', 'POST'])
def delete_index():
    return render_template('view1/delete_db_modal.html')

@app.route('/delete/db', methods=['GET', 'POST']) # note this does not currently delete anything from s3 bucket, only deletes from database
def delete_db():
    session = get_session()
    try:
        session.query(Text).delete()
        session.query(Audio).delete()
        session.query(Photo).delete()
        session.query(Video).delete()
        session.query(Powerpoint).delete()
        session.query(Project).delete()
        session.query(ProjectTextAssociation).delete()
        session.commit()
        session.close()
    except:
        session.rollback()
        raise
    print('database deleted')
    data = {'status': '200', 'message': 'Database successfully deleted'}
    formatted_data = json.dumps(data, indent=4).lstrip()
    return render_template('view1/deleted_modal.html', data=formatted_data)

@app.route('/view4', methods=['GET', 'POST'])
def create_powerpoint():
    session = get_session()
    video_data = session.query(Video).all()
    session.close()
    return render_template(
        'view4/view4.html', 
        data=video_data
    )

@app.route('/view4/submit', methods=['GET', 'POST'])
def create_powerpoint_submit():
    slide_title = request.form['slide_title']
    slide_body_text = request.form['slide_body']
    slide_video = request.form['video_select']
    # requires_formatting = request.form['requires_formatting']

    print(f'slide_title: {slide_title}, slide_body_text: {slide_body_text}, slide_video: {slide_video}')

    # if requires_formatting == 'true':
    normalized_text = slide_body_text.replace('\r\n', '\n').replace('\r', '\n')
    normalized_text = html.unescape(normalized_text)
    normalized_text = normalized_text.strip('"')


    # else:
    #     normalized_text = slide_body_text

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Using a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set slide title
    title = slide.shapes.title
    title.text = slide_title

    # Add body text
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(4))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = normalized_text


    print('slide_video FROM SERVER:', slide_video)

    # Add video
    ## first download video from s3
    video_url = slide_video
    video_name = video_url.split('/')[-1]
    print('video_name:', video_name)
    ## download video to /temp_outputs
    s3_client.download_file(BUCKET_NAME, video_name, f"./temp_outputs/{video_name}")
    video_path = f"./temp_outputs/{video_name}"
    image_path = video_path + '.jpg'

    def get_video_dimensions(video_path):
        """Get the width and height of the video."""
        cmd = [
            'ffprobe', 
            '-v', 'error',
            '-select_streams', 'v:0',
            '-show_entries', 'stream=width,height',
            '-of', 'json',
            video_path
        ]
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        if result.returncode != 0:
            raise Exception("ffprobe error: " + result.stderr)
        info = json.loads(result.stdout)
        width = info['streams'][0]['width']
        height = info['streams'][0]['height']
        return width, height
    
    # Get the dimensions of the video
    video_width, video_height = get_video_dimensions(video_path)

    ppt_height = Inches(1.5)
    aspect_ratio = video_width / video_height
    ppt_width = ppt_height * aspect_ratio

    # Positioning the video in the lower right corner of the slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    left = slide_width - ppt_width - Inches(0.5)  # 0.5 inches from the right side
    top = slide_height - ppt_height - Inches(0.5)  # 0.5 inches from the bottom

    # Use ffmpeg to extract a frame from the video
    subprocess.run([
        'ffmpeg',
        '-i', video_path,
        '-ss', '00:00:00',
        '-frames:v', '1',
        image_path
    ])

    slide.shapes.add_movie(
        video_path, left, top, ppt_width, ppt_height, 
        poster_frame_image=image_path, mime_type='video/mp4'
    )

    ## pptx name should be title of slide with date/time stamp
    now = datetime.datetime.now()
    date_time = now.strftime("%m-%d-%Y_%H-%M-%S")
    pptx_name = f'{slide_title}_{date_time}.pptx'

    ## first save to local
    prs.save(f'temp_outputs/{pptx_name}')

    ## save to s3
    s3_client.upload_file(f'temp_outputs/{pptx_name}', BUCKET_NAME, pptx_name)

    ## save to database in Powerpoint table
    session = get_session()
    try:
        new_powerpoint = Powerpoint(
            slide_title=slide_title,
            slide_body_text=slide_body_text,
            slide_video_url=slide_video,
            powerpoint_url=f'https://{BUCKET_NAME}.s3.amazonaws.com/{pptx_name}'
        )
        session.add(new_powerpoint)
        session.commit()
        session.close()
    except:
        session.rollback()
        raise

    ## then perform clean up

    ## delete any file that ends with .pptx, jpg, or mp4 in ./temp_outputs
    os.system("rm -rf temp_outputs/*.pptx")
    print("Deleted all .pptx files in temp_outputs")
    os.system("rm -rf temp_outputs/*.jpg")
    print("Deleted all .jpg files in temp_outputs")
    os.system("rm -rf temp_outputs/*.mp4")
    print("Deleted all .mp4 files in temp_outputs")

    return f'https://{BUCKET_NAME}.s3.amazonaws.com/{pptx_name}'


## endpoint for combining attempt 2 
@app.route('/view4/combine_slides', methods=['POST'])
def combine_slides():

    try:
        slides_data = request.get_json()
        print('Data object received in dictionary form: ', slides_data)

    except Exception as e:
        print(f"SERVER Error: {e}")
        return None
    
    print('Parsed Data:', slides_data)
    print('Data Type:', type(slides_data))

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Using a blank slide layout

    # Ensure that slides_data is a list
    if not isinstance(slides_data, list):
        raise ValueError("Data is not a list")

    for slide_data in slides_data:

        print('Working on PPT slide: ', slide_data['title'])

        slide_title = slide_data['title']
        slide_body_text = slide_data['body']
        slide_video = slide_data['videoUrl']

        normalized_text = slide_body_text.replace('\r\n', '\n').replace('\r', '\n')
        normalized_text = html.unescape(normalized_text)
        normalized_text = normalized_text.strip('"')

        # Create a new slide
        slide = prs.slides.add_slide(slide_layout)

        # Set slide title
        title = slide.shapes.title
        title.text = slide_title

        # Add body text
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(4))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = normalized_text

        # Add video
        ## first download video from s3
        video_url = slide_video
        video_name = video_url.split('/')[-1]
        print('video_name:', video_name)
        ## download video to /temp_outputs
        s3_client.download_file(BUCKET_NAME, video_name, f"./temp_outputs/{video_name}")
        video_path = f"./temp_outputs/{video_name}"
        image_path = video_path + '.jpg'

        def get_video_dimensions(video_path):
            """Get the width and height of the video."""
            cmd = [
                'ffprobe', 
                '-v', 'error',
                '-select_streams', 'v:0',
                '-show_entries', 'stream=width,height',
                '-of', 'json',
                video_path
            ]
            result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            if result.returncode != 0:
                raise Exception("ffprobe error: " + result.stderr)
            info = json.loads(result.stdout)
            width = info['streams'][0]['width']
            height = info['streams'][0]['height']
            return width, height
        
        # Get the dimensions of the video
        video_width, video_height = get_video_dimensions(video_path)

        ppt_height = Inches(1.5)
        aspect_ratio = video_width / video_height
        ppt_width = ppt_height * aspect_ratio

        # Positioning the video in the lower right corner of the slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        left = slide_width - ppt_width - Inches(0.5)  # 0.5 inches from the right side
        top = slide_height - ppt_height - Inches(0.5)  # 0.5 inches from the bottom

        # Use ffmpeg to extract a frame from the video
        subprocess.run([
            'ffmpeg',
            '-i', video_path,
            '-ss', '00:00:00',
            '-frames:v', '1',
            image_path
        ])

        slide.shapes.add_movie(
            video_path, left, top, ppt_width, ppt_height, 
            poster_frame_image=image_path, mime_type='video/mp4'
        )

        print('Finished on PPT slide: ', slide_data['title'])


    ## pptx name should be title of slide with date/time stamp
    now = datetime.datetime.now()
    date_time = now.strftime("%m-%d-%Y_%H-%M-%S")
    pptx_name = f'combined_test_{date_time}.pptx'

    ## first save to local
    print('Saving local version of combined...')
    prs.save(f'temp_outputs/{pptx_name}')

    ## save to s3
    print('Uploading to S3 version of combined...')
    s3_client.upload_file(f'temp_outputs/{pptx_name}', BUCKET_NAME, pptx_name)

    ## save to database in Powerpoint table
    session = get_session()
    try:
        new_powerpoint = Powerpoint(
            slide_title=slide_title,
            slide_body_text=slide_body_text,
            slide_video_url=slide_video,
            powerpoint_url=f'https://{BUCKET_NAME}.s3.amazonaws.com/{pptx_name}'
        )
        session.add(new_powerpoint)
        session.commit()
        session.close()
    except:
        session.rollback()
        raise

    ## then perform clean up
    ## delete any file that ends with .pptx, jpg, or mp4 in ./temp_outputs
    os.system("rm -rf temp_outputs/*.pptx")
    print("Deleted all .pptx files in temp_outputs")
    os.system("rm -rf temp_outputs/*.jpg")
    print("Deleted all .jpg files in temp_outputs")
    os.system("rm -rf temp_outputs/*.mp4")
    print("Deleted all .mp4 files in temp_outputs")

    jsonReturn = {'finalUrl': f'https://{BUCKET_NAME}.s3.amazonaws.com/{pptx_name}'}

    return jsonReturn

@app.route('/view5')
def view5():
    session = get_session()
    powerpoint_data = session.query(Powerpoint).all()
    session.close()
    return render_template(
        'view5/view5.html', 
        data=powerpoint_data
    )

if __name__ == '__main__':
    app.run(
        host='localhost',
        debug=True, 
        port=5005
    )
